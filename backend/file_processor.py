"""
File processing utilities for slides, notes, and YouTube transcripts
"""
import PyPDF2
from docx import Document
from pptx import Presentation
from youtube_transcript_api import YouTubeTranscriptApi
import re
from typing import Optional

class FileProcessor:
    @staticmethod
    def extract_pdf_text(file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            text = []
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text.append(page.extract_text())
            return "\n".join(text)
        except Exception as e:
            raise Exception(f"Error extracting PDF: {str(e)}")
    
    @staticmethod
    def extract_docx_text(file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            return "\n".join(text)
        except Exception as e:
            raise Exception(f"Error extracting DOCX: {str(e)}")
    
    @staticmethod
    def extract_pptx_text(file_path: str) -> str:
        """Extract text from PowerPoint file"""
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            raise Exception(f"Error extracting PPTX: {str(e)}")
    
    @staticmethod
    def extract_txt_text(file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            raise Exception(f"Error reading TXT: {str(e)}")
    
    @staticmethod
    def extract_youtube_video_id(url: str) -> Optional[str]:
        """Extract video ID from YouTube URL"""
        patterns = [
            r'(?:youtube\.com\/watch\?v=|youtu\.be\/)([a-zA-Z0-9_-]{11})',
            r'youtube\.com\/embed\/([a-zA-Z0-9_-]{11})',
        ]
        
        for pattern in patterns:
            match = re.search(pattern, url)
            if match:
                return match.group(1)
        return None
    
    @staticmethod
    def extract_playlist_id(url: str) -> Optional[str]:
        """Extract playlist ID from YouTube URL"""
        pattern = r'[?&]list=([a-zA-Z0-9_-]+)'
        match = re.search(pattern, url)
        return match.group(1) if match else None
    
    @staticmethod
    async def get_youtube_transcript(url: str) -> str:
        """Get transcript from YouTube video or playlist"""
        try:
            # Check if it's a playlist
            playlist_id = FileProcessor.extract_playlist_id(url)
            if playlist_id:
                return f"Playlist processing not fully implemented yet. Playlist ID: {playlist_id}\nNote: Extract individual video IDs from playlist and process each."
            
            # Single video
            video_id = FileProcessor.extract_youtube_video_id(url)
            if not video_id:
                raise ValueError("Invalid YouTube URL")
            
            # Get transcript
            transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
            transcript_text = " ".join([entry['text'] for entry in transcript_list])
            
            return transcript_text
        except Exception as e:
            raise Exception(f"Error getting YouTube transcript: {str(e)}")
    
    @staticmethod
    def process_file(file_path: str, file_type: str) -> str:
        """Process file based on type and return text content"""
        if file_type == "pdf":
            return FileProcessor.extract_pdf_text(file_path)
        elif file_type == "docx":
            return FileProcessor.extract_docx_text(file_path)
        elif file_type == "pptx":
            return FileProcessor.extract_pptx_text(file_path)
        elif file_type == "txt":
            return FileProcessor.extract_txt_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
